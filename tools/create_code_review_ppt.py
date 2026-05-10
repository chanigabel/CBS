from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt


OUT = Path("docs") / "excel_standardization_code_review_he.pptx"


SLIDES = [
    {
        "title": "Excel Standardization System - סקירת קוד",
        "bullets": [
            "מערכת סטנדרטיזציה לקבצי Excel של דיווחי מוסדות",
            "מיקוד: זרימה פעילה, שכבות, קבצים חשובים וסיכונים",
            "קהל יעד: ראש צוות וסקירת קוד",
            "מטרה: להבין איפה נכון לבדוק, לשפר ולהקשיח",
        ],
        "notes": "פתיחה קצרה: להסביר שהמצגת מתמקדת בנתיב הפעיל של המערכת ולא בכל ההיסטוריה של הקוד. להדגיש שהמטרה היא לא רק לתאר מה עובד, אלא לזהות נקודות בדיקה לקראת המשך תחזוקה.",
        "file": "README.md",
    },
    {
        "title": "מה המערכת עושה",
        "bullets": [
            "מקבלת Workbook מסוג xlsx או xlsm",
            "מחלצת גיליונות ושורות למבנה Dataset פנימי",
            "מנרמלת שמות, מגדר, תאריכים ומזהים",
            "מייצרת קובץ יצוא נקי בסכמה קבועה",
        ],
        "notes": "להסביר שהמערכת לא עורכת את הקובץ המקורי ישירות. העבודה נעשית על עותק עבודה ועל מבני נתונים פנימיים, ואז נכתב קובץ יצוא חדש.",
        "file": "README.md",
    },
    {
        "title": "הזרימה הפעילה המרכזית",
        "bullets": [
            "Upload: העלאת קובץ ויצירת session",
            "Extraction: זיהוי גיליונות, כותרות ושדות",
            "Standardization: יצירת שדות corrected",
            "Validation + Export: בדיקות שורה וכתיבת קובץ",
        ],
        "notes": "לעבור על השרשרת מקצה לקצה: upload, extraction, standardization, validation, export. חשוב להדגיש שהזרימה הפעילה מבוססת Dataset ולא מעבדי Excel ישנים.",
        "file": "webapp/services/standardization_service.py",
    },
    {
        "title": "שכבות ארכיטקטורה",
        "bullets": [
            "Web/API: endpoints, session וניהול קבצים",
            "I/O Layer: קריאה מ־Excel וזיהוי מבנה",
            "Processing: pipeline שמחבר את המנועים",
            "Engines: לוגיקה עסקית נקייה מתלות ב־Excel",
            "Validation/Export: סטטוסים וקובץ תוצאה",
        ],
        "notes": "להדגיש הפרדת אחריות: שכבת ה־I/O לא אמורה להכיל חוקים עסקיים, והמנועים לא אמורים להכיר openpyxl או FastAPI. זו נקודת בדיקה טובה בסקירת קוד.",
        "file": "src/excel_standardization/workbook_json_flow.py",
    },
    {
        "title": "קבצי src חשובים ותפקידם",
        "bullets": [
            "data_types.py: WorkbookDataset ו־SheetDataset",
            "excel_to_json_extractor.py: חילוץ שורות ל־JSON",
            "standardization_pipeline.py: תזמור כללי של הנרמול",
            "engines/*.py: חוקים נקודתיים לפי תחום",
            "export/export_engine.py: יצוא בסכמת יעד",
        ],
        "notes": "להציג את הקבצים כנקודות כניסה לסקירה. אם הזמן מוגבל, לפתוח קודם את pipeline, אחר כך מנוע אחד לדוגמה, ואז export.",
        "file": "src/excel_standardization/processing/standardization_pipeline.py",
    },
    {
        "title": "Standardization Pipeline",
        "bullets": [
            "מקבל SheetDataset ומחזיר Dataset מנורמל",
            "שומר ערכים מקוריים ומוסיף שדות עם סיומת corrected",
            "מזהה תבניות ברמת גיליון לפני ריצה על שורות",
            "אוסף כשלים וסטטיסטיקות במטא־דאטה",
        ],
        "notes": "להסביר שה־pipeline הוא adapter בין מבנה ה־JSON לבין המנועים. הוא גם המקום שבו נולדים נתוני בקרה כמו success_rate ו־failed_rows.",
        "file": "src/excel_standardization/processing/standardization_pipeline.py",
    },
    {
        "title": "מנועים: שם, מגדר, תאריך, מזהה",
        "bullets": [
            "NameEngine: ניקוי שמות והסרת שם משפחה לפי תבנית",
            "GenderEngine: המרה ל־1/2 לפי ייצוגים שונים",
            "DateEngine: פירוק, זיהוי פורמט וחוקי תאריך",
            "IdentifierEngine: תעודת זהות ודרכון",
        ],
        "notes": "להציג מנוע אחד לעומק, ואז להסביר שהשאר פועלים לפי אותו רעיון: פונקציות עסקיות יחסית מבודדות, עם תוצאה שה־pipeline כותב חזרה לשורה.",
        "file": "src/excel_standardization/engines/date_engine.py",
    },
    {
        "title": "שכבת Validation",
        "bullets": [
            "פועלת אחרי סטנדרטיזציה על ערכים corrected",
            "בודקת שדות חובה, כפילויות וטווחי תאריכים",
            "כותבת _validation_status ו־_validation_ok לכל שורה",
            "מופעלת רק לגיליונות מוסד מוכרים",
        ],
        "notes": "להדגיש שה־validation הוא post-processing. הוא לא מחליף את מנועי הנרמול, אלא מוסיף סטטוס עסקי לשורה אחרי שכבר יש ערכים מתוקנים.",
        "file": "src/excel_standardization/validation/institution_report_validator.py",
    },
    {
        "title": "שכבת Export",
        "bullets": [
            "כותבת Workbook חדש להורדה",
            "משתמשת בסכמת יעד קבועה ושמות גיליון ידועים",
            "מעדיפה ערכים corrected כשקיימים",
            "מסננת/ממפה שורות לפי לוגיקת היצוא",
        ],
        "notes": "להראות שה־export הוא חוזה חיצוני של המערכת. לכן שינוי קטן במיפוי עמודות או בשמות שדות יכול לשבור צרכנים downstream.",
        "file": "webapp/services/export_service.py",
    },
    {
        "title": "קוד פעיל מול קוד Legacy",
        "bullets": [
            "הנתיב הפעיל: Web/Dataset pipeline",
            "ה־orchestrator חוסם נתיבי Excel ישירים ישנים",
            "שמות backward-compatible נשמרים אך לא מפעילים legacy",
            "צריך להיזהר מבדיקות או imports שמחיים נתיב ישן",
        ],
        "notes": "להסביר למה זה חשוב בסקירת קוד: אם יש שני נתיבים שמבטיחים אותה תוצאה, יש סיכון לסטייה. כאן הקוד מנסה לכפות נתיב פעיל יחיד.",
        "file": "src/excel_standardization/orchestrator.py",
    },
    {
        "title": "סטנדרטים לפיתוח",
        "bullets": [
            "לוגיקה עסקית במנועים, לא ב־API ולא ב־I/O",
            "שינויים צריכים להיות דטרמיניסטיים וניתנים לבדיקה",
            "לא לשנות ערך מקור; להוסיף corrected",
            "להוסיף tests סביב כל כלל עסקי חדש",
            "לשמור על שמות שדות וסכמת יצוא יציבים",
        ],
        "notes": "זה שקף שמגדיר ציפיות לסקירת קוד. כל PR צריך להיבדק מול השאלות: איפה החוק נמצא, האם יש בדיקה, והאם נשמר החוזה של השדות והיצוא.",
        "file": "tests/",
    },
    {
        "title": "סיכונים וחוסרים בבדיקות",
        "bullets": [
            "קלט Excel לא צפוי: כותרות, merged cells ונוסחאות",
            "פער בין תיקון אוטומטי לבין שגיאת validation",
            "תלות בשמות גיליונות מוכרים ובמיפוי שדות",
            "מקורות חיצוניים חסרים לחלק מהבדיקות העסקיות",
            "סיכון רגרסיה מול התנהגות VBA",
        ],
        "notes": "להציג את הסיכונים כנקודות פעולה: להרחיב fixture-ים של Excel, להוסיף בדיקות export, ולתעד מקרים שבהם אין מקור אמת חיצוני.",
        "file": "src/excel_standardization/validation/institution_report_validator.py",
    },
    {
        "title": "שאלות לסקירת קוד",
        "bullets": [
            "האם הנתיב הפעיל ברור ואין עקיפות legacy?",
            "האם כל כלל עסקי נמצא בשכבה הנכונה?",
            "האם כשלי נרמול ו־validation מוצגים למשתמש נכון?",
            "האם export מכסה את כל השדות הנדרשים?",
            "אילו מקרים צריכים בדיקות רגרסיה נוספות?",
        ],
        "notes": "להשתמש בשקף הזה כדי להוביל דיון. המטרה היא לצאת עם החלטות: אילו סיכונים מתקנים עכשיו, אילו מתעדים, ואילו בדיקות מוסיפים.",
        "file": "webapp/services/export_validation.py",
    },
    {
        "title": "סיכום",
        "bullets": [
            "המערכת בנויה סביב Dataset pipeline פעיל",
            "המנועים מבודדים את החוקים העסקיים המרכזיים",
            "validation ו־export הם החוזה מול המשתמש והמערכות הבאות",
            "מוקדי שיפור: בדיקות Excel, כיסוי validation ותיעוד legacy",
        ],
        "notes": "לסגור בהמלצה מעשית: להתחיל סקירה מהזרימה הפעילה, לאמת את הסכמה וה־validation, ואז לבחור מנוע אחד לבדיקת עומק של איכות החוקים והבדיקות.",
        "file": "src/excel_standardization/workbook_json_flow.py",
    },
]


def set_rtl(paragraph, font_size=None, bold=False, color=None):
    p_pr = paragraph._p.get_or_add_pPr()
    p_pr.set("rtl", "1")
    paragraph.alignment = PP_ALIGN.RIGHT
    for run in paragraph.runs:
        run.font.name = "Arial"
        run.font.bold = bold
        if font_size:
            run.font.size = Pt(font_size)
        if color:
            run.font.color.rgb = color
        r_pr = run._r.get_or_add_rPr()
        r_pr.set("rtl", "1")
        r_pr.set("lang", "he-IL")


def add_textbox(slide, left, top, width, height, text, size, bold=False, color=RGBColor(31, 41, 55)):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    set_rtl(p, size, bold, color)
    return box


def add_bullets(slide, bullets):
    box = slide.shapes.add_textbox(Cm(1.25), Cm(3.15), Cm(22.7), Cm(9.1))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(23)
        p.space_after = Pt(8)
        p.alignment = PP_ALIGN.RIGHT
        p_pr = p._p.get_or_add_pPr()
        p_pr.set("rtl", "1")
        p_pr.set("marR", "342900")
        p_pr.set("indent", "-228600")
        for run in p.runs:
            run.font.name = "Arial"
            run.font.color.rgb = RGBColor(39, 39, 42)
            r_pr = run._r.get_or_add_rPr()
            r_pr.set("rtl", "1")
            r_pr.set("lang", "he-IL")
    return box


def add_notes(slide, notes, file_rec):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    notes_tf.text = f"{notes}\n\nקובץ לפתיחה בלייב: {file_rec}"
    for p in notes_tf.paragraphs:
        set_rtl(p, 12, False, RGBColor(0, 0, 0))


def build():
    prs = Presentation()
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(14.288)
    blank = prs.slide_layouts[6]

    navy = RGBColor(15, 23, 42)
    accent = RGBColor(20, 184, 166)
    muted = RGBColor(100, 116, 139)
    bg = RGBColor(248, 250, 252)

    for i, data in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg

        # Top accent bar
        bar = slide.shapes.add_shape(1, Cm(0), Cm(0), prs.slide_width, Cm(0.22))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        add_textbox(slide, Cm(1.2), Cm(0.75), Cm(22.9), Cm(1.25), data["title"], 31 if i == 1 else 28, True, navy)
        add_bullets(slide, data["bullets"])

        # Live file recommendation footer
        footer = f"קובץ לפתיחה בלייב: {data['file']}"
        add_textbox(slide, Cm(1.2), Cm(12.55), Cm(22.9), Cm(0.62), footer, 13, False, muted)
        add_textbox(slide, Cm(1.2), Cm(13.2), Cm(22.9), Cm(0.42), f"{i}/14", 10, False, muted)

        add_notes(slide, data["notes"], data["file"])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT.resolve())


if __name__ == "__main__":
    build()
